"""
KISTi 프로젝트 기획 자료 추출 스크립트 v3
- 카테고리 7개 + Notion MD (2024 작업현황)
- PDF 텍스트 추출 (PyMuPDF)
- PPT → 슬라이드별 텍스트 + 이미지
- Excel → 시트별 테이블 데이터
- Notion MD → 마크다운 원문
"""
import os
import json
import re
import hashlib
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openpyxl import load_workbook
from PIL import Image
import io
import fitz  # PyMuPDF

# --- Config ---
BASE_DIR = Path(__file__).parent
DATA_DIR = BASE_DIR / "data"
IMAGES_DIR = DATA_DIR / "images"

SOURCE_DIRS = [
    Path(r"C:\Users\Lyu\Desktop\[2024]KISTi"),
    Path(r"C:\Users\Lyu\Desktop\[2025]KISTi"),
]

NOTION_DIR = BASE_DIR / "notion_md_2024"

# ============================================================
# 7 unified categories
# ============================================================
CATEGORY_RULES = [
    ("곤충잡기",   ["곤충", "insect", "곤충잡기", "곤충채집", "곤충배치", "곤충 생성", "스폰 범위"]),
    ("공놀이",     ["공놀이", "ball", "전광판"]),
    ("인지검사",   ["인지평가", "인지검사", "cognitive", "인지훈련"]),
    ("균형검사",   ["균형", "force plate", "balance"]),
    ("운동성검사", ["운동성", "운동 콘텐츠", "운동콘텐츠", "stability", "운동효과"]),
    ("런처",       ["런처", "launcher", "cctv", "교수자", "학습자", "수업관리", "디바이스 매칭"]),
]

def detect_category(filepath: Path) -> str:
    full_lower = str(filepath).lower()
    name_lower = filepath.name.lower()
    search_str = full_lower + " " + name_lower
    for cat, keywords in CATEGORY_RULES:
        for kw in keywords:
            if kw.lower() in search_str:
                return cat
    return "공통/기타"

# ============================================================
# Helpers
# ============================================================
def ensure_dirs():
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)

def extract_date_from_filename(filename: str) -> str:
    m = re.search(r'(20\d{2})(\d{2})(\d{2})', filename)
    if m:
        return f"{m.group(1)}-{m.group(2)}-{m.group(3)}"
    m = re.search(r'(\d{2})(\d{2})(\d{2})', filename)
    if m and 20 <= int(m.group(1)) <= 30:
        return f"20{m.group(1)}-{m.group(2)}-{m.group(3)}"
    return ""

def safe_image_name(filepath: Path, slide_idx: int, img_idx: int, ext="png") -> str:
    h = hashlib.md5(str(filepath).encode()).hexdigest()[:8]
    return f"{h}_s{slide_idx}_i{img_idx}.{ext}"

# ============================================================
# PPT Extraction
# ============================================================
def extract_pptx(filepath: Path) -> dict:
    try:
        prs = Presentation(str(filepath))
    except Exception as e:
        print(f"  [ERROR] {filepath.name}: {e}")
        return {"error": str(e), "slides": []}

    slides = []
    for si, slide in enumerate(prs.slides):
        info = {"index": si + 1, "texts": [], "images": [], "tables": []}

        for shape in slide.shapes:
            _extract_shape_content(shape, filepath, si, info)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for gs in shape.shapes:
                        _extract_shape_content(gs, filepath, si, info)
                except:
                    pass

        slides.append(info)

    return {"slide_count": len(prs.slides), "slides": slides}


def _extract_shape_content(shape, filepath, si, info):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                is_heading = any(r.font.size and r.font.size >= Pt(18) for r in para.runs)
                info["texts"].append({"content": text, "is_heading": is_heading})

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            blob = shape.image.blob
            ext = shape.image.content_type.split('/')[-1]
            if ext == 'jpeg': ext = 'jpg'
            name = safe_image_name(filepath, si, len(info["images"]), ext)
            (IMAGES_DIR / name).write_bytes(blob)
            try:
                w, h = Image.open(io.BytesIO(blob)).size
            except:
                w, h = 0, 0
            info["images"].append({"filename": name, "width": w, "height": h})
        except:
            pass

    if shape.has_table:
        rows = []
        for row in shape.table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        info["tables"].append(rows)

# ============================================================
# Excel Extraction
# ============================================================
def extract_xlsx(filepath: Path) -> dict:
    try:
        wb = load_workbook(str(filepath), data_only=True)
    except Exception as e:
        print(f"  [ERROR] {filepath.name}: {e}")
        return {"error": str(e), "sheets": []}

    sheets = []
    for sname in wb.sheetnames:
        ws = wb[sname]
        merged_map = {}
        for mr in ws.merged_cells.ranges:
            for r in range(mr.min_row, mr.max_row + 1):
                for c in range(mr.min_col, mr.max_col + 1):
                    if not (r == mr.min_row and c == mr.min_col):
                        merged_map[(r, c)] = (mr.min_row, mr.min_col)

        rows = []
        for ri in range(1, min((ws.max_row or 1) + 1, 200)):
            row = []
            for ci in range(1, min((ws.max_column or 1) + 1, 30)):
                src = merged_map.get((ri, ci), (ri, ci))
                val = ws.cell(row=src[0], column=src[1]).value
                if val is None:
                    val = ""
                elif isinstance(val, datetime):
                    val = val.strftime("%Y-%m-%d")
                else:
                    val = str(val)
                row.append(val)
            if any(v.strip() for v in row):
                rows.append(row)

        if rows:
            sheets.append({"name": sname, "rows": rows, "col_count": ws.max_column, "row_count": len(rows)})
    wb.close()
    return {"sheets": sheets}

# ============================================================
# PDF Extraction (PyMuPDF)
# ============================================================
def extract_pdf(filepath: Path) -> dict:
    try:
        doc = fitz.open(str(filepath))
    except Exception as e:
        print(f"  [ERROR] {filepath.name}: {e}")
        return {"error": str(e), "pages": []}

    pages = []
    for pi in range(len(doc)):
        page = doc[pi]
        text = page.get_text("text").strip()
        
        # Extract images from PDF
        images = []
        for img_idx, img in enumerate(page.get_images(full=True)):
            try:
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                if pix.n - pix.alpha > 3:  # CMYK
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                name = safe_image_name(filepath, pi, img_idx, "png")
                pix.save(str(IMAGES_DIR / name))
                images.append({"filename": name, "width": pix.width, "height": pix.height})
                pix = None
            except:
                pass
        
        if text or images:
            pages.append({"index": pi + 1, "text": text, "images": images})

    doc.close()
    return {"page_count": len(pages), "pages": pages}

# ============================================================
# Markdown Extraction (Notion Export)
# ============================================================
def extract_md(filepath: Path) -> dict:
    try:
        text = filepath.read_text(encoding='utf-8')
    except Exception as e:
        print(f"  [ERROR] {filepath.name}: {e}")
        return {"error": str(e), "text": ""}
    return {"text": text, "char_count": len(text)}

def clean_notion_title(filename: str) -> str:
    """Remove Notion UUID hash from filename."""
    name = re.sub(r'\s+[0-9a-f]{32}\.md$', '.md', filename)
    name = re.sub(r'\s+[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}\.md$', '.md', name)
    return name.replace('.md', '')

# ============================================================
# Main
# ============================================================
def process_all():
    ensure_dirs()
    all_docs = []
    fid = 0

    # --- Process PPT/Excel/PDF from source directories ---
    for src_dir in SOURCE_DIRS:
        if not src_dir.exists():
            continue
        year = "2024" if "2024" in src_dir.name else "2025"
        print(f"\n{'='*60}\n{src_dir.name}\n{'='*60}")

        for fp in sorted(src_dir.rglob("*")):
            if fp.is_dir() or fp.name.startswith("~$"):
                continue
            ext = fp.suffix.lower()
            if ext not in ('.pptx', '.xlsx', '.pdf'):
                continue

            fid += 1
            cat = detect_category(fp)
            date = extract_date_from_filename(fp.name)

            doc = {
                "id": fid,
                "filename": fp.name,
                "relative_path": str(fp.relative_to(src_dir)),
                "year": year,
                "category": cat,
                "date": date,
                "file_type": ext.lstrip('.'),
                "file_size": fp.stat().st_size,
            }

            print(f"[{fid:3d}] {cat:10s} | {ext:5s} | {fp.name}")

            if ext == '.pptx':
                doc["content"] = extract_pptx(fp)
            elif ext == '.xlsx':
                doc["content"] = extract_xlsx(fp)
            elif ext == '.pdf':
                doc["content"] = extract_pdf(fp)

            all_docs.append(doc)

    # --- Process Notion Markdown files ---
    if NOTION_DIR.exists():
        print(f"\n{'='*60}\nNotion Markdown (2024)\n{'='*60}")
        for fp in sorted(NOTION_DIR.rglob("*.md")):
            if fp.is_dir():
                continue
            fid += 1
            
            # Determine subfolder for LNB tree
            try:
                rel = fp.relative_to(NOTION_DIR)
            except:
                rel = Path(fp.name)
            
            parts = list(rel.parts)
            # Remove the top-level Notion folder structure to simplify
            # "개인 페이지 & 공유된 페이지/2024 KISTi 공통 작업 진행 현황/..."
            # We want to keep from "2024 KISTi..." onwards
            clean_parts = []
            for p in parts:
                if '2024 KISTi' in p or len(clean_parts) > 0:
                    clean_parts.append(p)
            if not clean_parts:
                clean_parts = parts
            
            # Determine sub-category within Notion
            notion_folder = clean_parts[1] if len(clean_parts) > 2 else ""
            
            title = clean_notion_title(fp.name)
            date = extract_date_from_filename(fp.name)
            
            doc = {
                "id": fid,
                "filename": fp.name,
                "title": title,
                "relative_path": str(rel),
                "notion_folder": notion_folder,
                "year": "2024",
                "category": "2024 작업현황",
                "date": date,
                "file_type": "md",
                "file_size": fp.stat().st_size,
                "content": extract_md(fp),
            }

            print(f"[{fid:3d}] {'Notion':10s} | .md   | {title}")
            all_docs.append(doc)

    # --- Summary per category ---
    cat_summary = {}
    for d in all_docs:
        c = d["category"]
        if c not in cat_summary:
            cat_summary[c] = {"count": 0, "pptx": 0, "xlsx": 0, "pdf": 0, "md": 0}
        cat_summary[c]["count"] += 1
        ft = d["file_type"]
        if ft in cat_summary[c]:
            cat_summary[c][ft] += 1

    output = DATA_DIR / "content.json"
    with open(output, 'w', encoding='utf-8') as f:
        json.dump({
            "generated_at": datetime.now().isoformat(),
            "total_documents": len(all_docs),
            "category_summary": cat_summary,
            "documents": all_docs,
        }, f, ensure_ascii=False, indent=2)

    print(f"\n{'='*60}")
    print(f"✅ {len(all_docs)} documents → {output}")
    for cat, s in sorted(cat_summary.items()):
        print(f"  {cat}: {s['count']}개 (PPT:{s['pptx']}, Excel:{s['xlsx']}, PDF:{s['pdf']}, MD:{s.get('md',0)})")

if __name__ == "__main__":
    process_all()
